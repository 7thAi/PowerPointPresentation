import os
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import filedialog, messagebox, ttk

def count_photos(folder_path):
    """Считает количество изображений в указанной папке."""
    count = sum(1 for file in os.listdir(folder_path) if file.endswith(('jpg', 'jpeg', 'png', 'gif', 'heic')))
    return count

def create_presentation(folder_path, template_path, positions, output_name):
    """Создает презентацию PowerPoint, используя шаблон и добавляя изображения."""
    photo_count = count_photos(folder_path)
    num_slides = -(-photo_count // 6)  # Округление вверх
    prs = Presentation(template_path)

    for _ in range(num_slides - 1):
        prs.slides.add_slide(prs.slides[1].slide_layout)  # Клонируем 2-й слайд

    for i in range(num_slides):
        slide = prs.slides[i + 1]  # Начинаем со второго слайда

        for j in range(6):
            if (i * 6 + j) < photo_count:
                photo_file = os.listdir(folder_path)[i * 6 + j]
                img_path = os.path.join(folder_path, photo_file)

                # Указываем позиции для каждой фотографии
                left_position, top_position = positions[j]
                left = Inches(left_position)
                top = Inches(top_position)

                slide.shapes.add_picture(img_path, left, top, width=Inches(9.74 / 2.54), height=Inches(7.09 / 2.54))

                # Добавляем текстовое поле с именем файла под изображением
                textbox_left = left
                textbox_top = top + Inches(7.09 / 2.54 + 0.5)
                textbox_width = Inches(9.74 / 2.54)
                textbox_height = Inches(0.5)

                textbox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
                text_frame = textbox.text_frame
                text_frame.text = photo_file

    prs.save(output_name)

def browse_folder(entry):
    folder_path = filedialog.askdirectory()
    if folder_path:
        entry.delete(0, tk.END)
        entry.insert(0, folder_path)

def browse_template(entry):
    template_path = filedialog.askopenfilename(title="Выберите шаблон презентации", filetypes=[("PowerPoint files", "*.pptx")])
    if template_path:
        entry.delete(0, tk.END)
        entry.insert(0, template_path)

def save_presentation():
    template_path = template_entry.get()
    folder_path = folder_entry.get()
    output_name = output_entry.get()

    if not template_path or not folder_path or not output_name:
        messagebox.showerror("Ошибка", "Пожалуйста, заполните все поля.")
        return

    # Добавляем .pptx, если пользователь не указал его
    if not output_name.endswith(".pptx"):
        output_name += ".pptx"

    # Позиции для фотографий
    positions = [
        (2.42 / 2.54, 4.42 / 2.54),
        (12.86 / 2.54, 4.42 / 2.54),
        (23.32 / 2.54, 4.42 / 2.54),
        (2.48 / 2.54, 14.53 / 2.54),
        (12.86 / 2.54, 14.53 / 2.54),
        (23.32 / 2.54, 14.53 / 2.54),
    ]

    create_presentation(folder_path, template_path, positions, output_name)
    root.destroy()  # Закрыть интерфейс после успешного создания

# Создаем основной интерфейс
root = tk.Tk()
root.title("Создание презентации")
root.geometry("500x200")  # Установить размеры окна
root.configure(bg="#F0F0F0")  # Цвет фона

frame = ttk.Frame(root, padding="10")
frame.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)

# Поле для выбора шаблона презентации
template_label = ttk.Label(frame, text="Шаблон презентации:", background="#F0F0F0")
template_label.grid(row=0, column=0, sticky=tk.E)

template_entry = ttk.Entry(frame, width=40)
template_entry.grid(row=0, column=1, padx=5)

template_button = ttk.Button(frame, text="Обзор...", command=lambda: browse_template(template_entry))
template_button.grid(row=0, column=2)

# Поле для выбора папки с фото
folder_label = ttk.Label(frame, text="Папка с фото:", background="#F0F0F0")
folder_label.grid(row=1, column=0, sticky=tk.E)

folder_entry = ttk.Entry(frame, width=40)
folder_entry.grid(row=1, column=1, padx=5)

browse_button = ttk.Button(frame, text="Обзор...", command=lambda: browse_folder(folder_entry))
browse_button.grid(row=1, column=2)

# Поле для имени файла
output_label = ttk.Label(frame, text="Имя файла:", background="#F0F0F0")
output_label.grid(row=2, column=0, sticky=tk.E)

output_entry = ttk.Entry(frame, width=40)
output_entry.grid(row=2, column=1, padx=5)

save_button = ttk.Button(frame, text="Создать презентацию", command=save_presentation)
save_button.grid(row=3, column=1, columnspan=1, pady=10)

root.mainloop()